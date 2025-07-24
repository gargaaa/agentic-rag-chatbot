from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd
import os

def parse_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
    elif ext == ".docx":
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    elif ext == ".pptx":
        prs = Presentation(file_path)
        return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text)
    elif ext == ".csv":
        df = pd.read_csv(file_path)
        return df.to_string()
    elif ext in [".txt", ".md"]:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    return ""
